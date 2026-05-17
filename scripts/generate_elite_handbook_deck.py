#!/usr/bin/env python3
"""Generate a dense, polished slideshow from the appointment setter handbook."""

from __future__ import annotations

import argparse
import math
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass
class Section:
    level: int
    title: str
    lines: list[str] = field(default_factory=list)


PALETTE = {
    "bg": RGBColor(8, 15, 30),
    "bg_alt": RGBColor(12, 23, 44),
    "panel": RGBColor(19, 33, 58),
    "line": RGBColor(45, 74, 120),
    "accent": RGBColor(230, 184, 92),
    "text": RGBColor(242, 247, 255),
    "muted": RGBColor(164, 180, 210),
}


def clean_inline_markdown(text: str) -> str:
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = text.replace("**", "").replace("__", "").replace("`", "")
    return text.strip()


def parse_markdown(markdown_text: str) -> tuple[str, list[Section]]:
    title = "Appointment Setter Handbook"
    sections: list[Section] = []
    current: Section | None = None

    for raw_line in markdown_text.splitlines():
        heading_match = re.match(r"^(#{1,6})\s+(.*)$", raw_line)
        if heading_match:
            if current is not None:
                sections.append(current)
            level = len(heading_match.group(1))
            heading_text = clean_inline_markdown(heading_match.group(2))
            if level == 1:
                title = heading_text
                current = None
            else:
                current = Section(level=level, title=heading_text)
            continue

        if current is not None:
            current.lines.append(raw_line.rstrip())

    if current is not None:
        sections.append(current)

    return title, sections


def group_sections(sections: list[Section]) -> list[Section]:
    """Group nested headings under each level-2 section for compact slides."""
    grouped: list[Section] = []
    current: Section | None = None

    for section in sections:
        if section.level == 2:
            if current is not None:
                grouped.append(current)
            current = Section(level=2, title=section.title, lines=list(section.lines))
            continue

        if current is None:
            current = Section(level=2, title="Overview", lines=[])

        if current.lines and current.lines[-1].strip():
            current.lines.append("")
        current.lines.append(f"### {section.title}")
        current.lines.extend(section.lines)

    if current is not None:
        grouped.append(current)

    return grouped


def normalize_lines(lines: Iterable[str]) -> list[str]:
    normalized: list[str] = []
    previous_empty = False

    for raw in lines:
        stripped = raw.strip()
        if stripped == "---":
            continue
        if not stripped:
            if not previous_empty:
                normalized.append("")
            previous_empty = True
            continue
        previous_empty = False
        normalized.append(clean_inline_markdown(stripped))

    while normalized and normalized[0] == "":
        normalized.pop(0)
    while normalized and normalized[-1] == "":
        normalized.pop()
    return normalized


def chunk_lines(lines: list[str], max_lines: int = 20, max_chars: int = 1700) -> list[list[str]]:
    if not lines:
        return []

    chunks: list[list[str]] = []
    current: list[str] = []
    visual_line_total = 0.0
    char_total = 0

    for line in lines:
        line_weight = 1.25 if line.startswith("### ") else 1.0
        next_line_total = visual_line_total + line_weight
        next_char_total = char_total + len(line)
        overflow = current and (
            next_line_total > max_lines or next_char_total > max_chars
        )

        if overflow:
            chunks.append(current)
            current = []
            visual_line_total = 0.0
            char_total = 0

        current.append(line)
        visual_line_total += line_weight
        char_total += len(line)

    if current:
        chunks.append(current)
    return chunks


def line_weight(line: str) -> float:
    if not line:
        return 0.45
    if line.startswith("### "):
        return 1.35
    return 1.0


def rebalance_chunks(chunks: list[list[str]], max_chunks: int = 2) -> list[list[str]]:
    """Reduce pagination count while preserving line order."""
    if len(chunks) <= max_chunks:
        return chunks

    lines: list[str] = [line for chunk in chunks for line in chunk]
    total_weight = sum(line_weight(line) for line in lines)
    if total_weight == 0:
        return [lines]

    target = total_weight / max_chunks
    rebalanced: list[list[str]] = []
    current: list[str] = []
    running = 0.0
    chunk_index = 1

    for idx, line in enumerate(lines):
        current.append(line)
        running += line_weight(line)

        lines_remaining = len(lines) - (idx + 1)
        chunks_remaining = max_chunks - chunk_index
        must_split = chunks_remaining > 0 and lines_remaining >= chunks_remaining
        near_boundary = running >= (target * chunk_index)
        split_point = line == "" or line.startswith("### ")

        if must_split and near_boundary and split_point:
            rebalanced.append(current)
            current = []
            chunk_index += 1

    if current:
        rebalanced.append(current)

    while len(rebalanced) > max_chunks:
        rebalanced[-2].extend(rebalanced[-1])
        rebalanced.pop()
    while len(rebalanced) < max_chunks:
        rebalanced.append([])

    return [chunk for chunk in rebalanced if chunk]


def apply_background(slide, slide_width, slide_height) -> None:
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALETTE["bg"]
    bg.line.fill.background()

    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_width, Inches(0.22)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PALETTE["accent"]
    overlay.line.fill.background()

    rail = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), Inches(0.7), Inches(0.08), Inches(5.8)
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = PALETTE["line"]
    rail.line.fill.background()


def add_footer(slide, slide_number: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.58), Inches(6.86), Inches(12.15), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Elite Appointment Setter Pitch Booklet  |  Slide {slide_number}"
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.color.rgb = PALETTE["muted"]
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs: Presentation, deck_title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, prs.slide_width, prs.slide_height)

    brand = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(10.5), Inches(0.45))
    brand_tf = brand.text_frame
    brand_tf.clear()
    bp = brand_tf.paragraphs[0]
    bp.text = "ELITE FIELD EXECUTION TRAINING"
    bp.font.name = "Aptos Display"
    bp.font.bold = True
    bp.font.size = Pt(16)
    bp.font.color.rgb = PALETTE["accent"]

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.28), Inches(11.7), Inches(2.4))
    tf = title_box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = deck_title
    p.font.name = "Aptos Display"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = PALETTE["text"]

    subtitle_box = slide.shapes.add_textbox(Inches(0.95), Inches(3.8), Inches(11.1), Inches(1.2))
    stf = subtitle_box.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.name = "Aptos"
    sp.font.size = Pt(20)
    sp.font.color.rgb = PALETTE["muted"]

    add_footer(slide, 1)


def add_agenda_slide(prs: Presentation, sections: list[Section], slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, prs.slide_width, prs.slide_height)

    heading_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.55), Inches(10.8), Inches(0.8))
    htf = heading_box.text_frame
    htf.clear()
    hp = htf.paragraphs[0]
    hp.text = "Training Roadmap"
    hp.font.name = "Aptos Display"
    hp.font.bold = True
    hp.font.size = Pt(31)
    hp.font.color.rgb = PALETTE["text"]

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.33), Inches(11.25), Inches(5.15)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALETTE["panel"]
    panel.line.color.rgb = PALETTE["line"]
    panel.line.width = Pt(1.2)

    agenda = slide.shapes.add_textbox(Inches(1.15), Inches(1.63), Inches(10.55), Inches(4.55))
    tf = agenda.text_frame
    tf.clear()
    tf.word_wrap = True

    max_items = 18
    visible_sections = sections[:max_items]
    col_count = 2
    per_col = math.ceil(len(visible_sections) / col_count)
    columns: list[list[str]] = []
    for col_idx in range(col_count):
        start = col_idx * per_col
        end = min(start + per_col, len(visible_sections))
        if start < len(visible_sections):
            columns.append(
                [f"{i + 1}. {sec.title}" for i, sec in enumerate(visible_sections[start:end], start=start)]
            )

    col_width = Inches(5.05)
    for col_idx, col_items in enumerate(columns):
        x = Inches(1.15) + Inches(5.23 * col_idx)
        box = slide.shapes.add_textbox(x, Inches(1.63), col_width, Inches(4.55))
        btf = box.text_frame
        btf.clear()
        for idx, item in enumerate(col_items):
            paragraph = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
            paragraph.text = item
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(15)
            paragraph.font.color.rgb = PALETTE["text"]
            paragraph.space_after = Pt(7)

    if len(sections) > max_items:
        note = slide.shapes.add_textbox(Inches(1.15), Inches(6.15), Inches(9.5), Inches(0.3))
        ntf = note.text_frame
        ntf.clear()
        np = ntf.paragraphs[0]
        np.text = "Deck continues with additional subsections in detailed flow."
        np.font.name = "Aptos"
        np.font.size = Pt(10)
        np.font.color.rgb = PALETTE["muted"]

    add_footer(slide, slide_number)


def estimated_units(lines: list[str], width_chars: int) -> float:
    units = 0.0
    for line in lines:
        if not line:
            units += 0.45
            continue

        text = line.replace("### ", "", 1)
        text = re.sub(r"^[-*]\s+", "", text)
        wraps = max(1, math.ceil(len(text) / width_chars))
        if line.startswith("### "):
            units += 1.0 + (0.75 * wraps)
        else:
            units += 0.8 * wraps
    return units


def _render_text_lines(text_frame, lines: list[str], body_font_size: int) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Inches(0.02)
    text_frame.margin_right = Inches(0.02)
    text_frame.margin_top = Inches(0.01)
    text_frame.margin_bottom = Inches(0.01)

    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.space_after = Pt(2)
        paragraph.line_spacing = 1.0
        paragraph.alignment = PP_ALIGN.LEFT

        if not line:
            paragraph.text = ""
            continue

        is_subheading = line.startswith("### ")
        if is_subheading:
            paragraph.text = line.replace("### ", "", 1)
            paragraph.font.name = "Aptos"
            paragraph.font.bold = True
            paragraph.font.size = Pt(body_font_size + 1)
            paragraph.font.color.rgb = PALETTE["accent"]
            paragraph.space_after = Pt(3)
            continue

        bullet_match = re.match(r"^[-*]\s+(.*)$", line)
        paragraph.text = f"• {bullet_match.group(1)}" if bullet_match else line
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(body_font_size)
        paragraph.font.color.rgb = PALETTE["text"]


def split_for_columns(content_chunk: list[str]) -> tuple[list[str], list[str] | None]:
    if len(content_chunk) <= 11:
        return content_chunk, None

    total_units = sum(line_weight(line) for line in content_chunk)
    target = total_units / 2
    running = 0.0
    split_index = len(content_chunk) // 2

    for idx, line in enumerate(content_chunk):
        running += line_weight(line)
        if running >= target and idx < len(content_chunk) - 1:
            split_index = idx + 1
            break

    window_start = max(1, split_index - 4)
    window_end = min(len(content_chunk) - 1, split_index + 4)
    for idx in range(window_start, window_end + 1):
        if content_chunk[idx - 1] == "" or content_chunk[idx - 1].startswith("### "):
            split_index = idx
            break

    return content_chunk[:split_index], content_chunk[split_index:]


def add_content_slide(
    prs: Presentation, heading: str, content_chunk: list[str], slide_number: int
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, prs.slide_width, prs.slide_height)

    heading_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.45), Inches(11.2), Inches(0.8))
    htf = heading_box.text_frame
    htf.clear()
    htf.vertical_anchor = MSO_ANCHOR.TOP
    hp = htf.paragraphs[0]
    hp.text = heading
    hp.font.name = "Aptos Display"
    hp.font.bold = True
    title_size = 27
    if len(heading) > 70:
        title_size = 23
    elif len(heading) > 52:
        title_size = 25
    hp.font.size = Pt(title_size)
    hp.font.color.rgb = PALETTE["text"]

    body_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.2), Inches(11.2), Inches(5.45)
    )
    body_panel.fill.solid()
    body_panel.fill.fore_color.rgb = PALETTE["panel"]
    body_panel.line.color.rgb = PALETTE["line"]
    body_panel.line.width = Pt(1.25)

    left_chunk, right_chunk = split_for_columns(content_chunk)
    if right_chunk is None:
        body = slide.shapes.add_textbox(Inches(1.08), Inches(1.42), Inches(10.8), Inches(5.0))
        units = estimated_units(left_chunk, width_chars=95)
        font_size = 16
        if units > 28:
            font_size = 14
        elif units > 23:
            font_size = 15
        _render_text_lines(body.text_frame, left_chunk, body_font_size=font_size)
    else:
        left_box = slide.shapes.add_textbox(Inches(1.05), Inches(1.4), Inches(5.2), Inches(5.0))
        right_box = slide.shapes.add_textbox(Inches(6.4), Inches(1.4), Inches(5.2), Inches(5.0))

        left_units = estimated_units(left_chunk, width_chars=44)
        right_units = estimated_units(right_chunk, width_chars=44)
        max_units = max(left_units, right_units)
        font_size = 15
        if max_units > 30:
            font_size = 12
        elif max_units > 26:
            font_size = 13
        elif max_units > 22:
            font_size = 14

        _render_text_lines(left_box.text_frame, left_chunk, body_font_size=font_size)
        _render_text_lines(right_box.text_frame, right_chunk, body_font_size=font_size)

    add_footer(slide, slide_number)


def build_deck(markdown_path: Path, output_path: Path) -> int:
    markdown = markdown_path.read_text(encoding="utf-8")
    deck_title, sections = parse_markdown(markdown)
    grouped_sections = group_sections(sections)

    subtitle = (
        "Door-to-Door Savings and Net Metering Pitch Execution Guide"
    )
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_number = 1
    add_title_slide(prs, deck_title, subtitle)

    slide_number += 1
    add_agenda_slide(prs, grouped_sections, slide_number)

    for section in grouped_sections:
        normalized_lines = normalize_lines(section.lines)
        chunks = chunk_lines(normalized_lines)
        chunks = rebalance_chunks(chunks, max_chunks=2)

        if not chunks:
            continue

        for part_index, chunk in enumerate(chunks, start=1):
            title = section.title
            if len(chunks) > 1:
                title = f"{section.title} ({part_index}/{len(chunks)})"
            slide_number += 1
            add_content_slide(prs, title, chunk, slide_number)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return slide_number


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate a polished PPTX deck from handbook markdown."
    )
    parser.add_argument(
        "--input",
        default="APPOINTMENT_SETTER_HANDBOOK.md",
        help="Path to source markdown file.",
    )
    parser.add_argument(
        "--output",
        default="deliverables/APPOINTMENT_SETTER_HANDBOOK_ELITE_DECK.pptx",
        help="Path for generated PPTX output.",
    )
    args = parser.parse_args()

    slide_count = build_deck(Path(args.input), Path(args.output))
    print(f"Created {args.output} with {slide_count} slides.")


if __name__ == "__main__":
    main()
